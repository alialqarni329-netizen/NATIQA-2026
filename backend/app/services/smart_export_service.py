"""
Smart Export Service — AI-powered document analysis + multi-format generation
═══════════════════════════════════════════════════════════════════════════════
Supported output formats:
  • Excel  (.xlsx)  — pandas + openpyxl
  • Word   (.docx)  — python-docx
  • PDF    (.pdf)   — reportlab (Arabic-aware via arabic-reshaper + bidi)
  • PowerPoint (.pptx) — python-pptx
  • Power BI  (.json) — structured dataset JSON (importable to Power BI Desktop)

Flow:
  1. extract_full_text(bytes, filename)  — full text extraction
  2. ai_analyze(text, format, export_type, llm)  — LLM returns structured JSON
  3. generate_*(ai_result) — one generator per format
  4. Returns (BytesIO, mime_type, filename)
"""

from __future__ import annotations

import io
import json
import re
import textwrap
from datetime import datetime, timezone
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

import structlog

log = structlog.get_logger()

# ── Optional imports (graceful fallback) ──────────────────────────────────────
try:
    import pandas as pd
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.chart import BarChart, Reference
    _HAS_EXCEL = True
except ImportError:
    _HAS_EXCEL = False

try:
    from docx import Document as DocxDocument
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PI, Pt as PPt, Emu
    from pptx.dml.color import RGBColor as PptRGB
    from pptx.enum.text import PP_ALIGN
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import arabic_reshaper
    from bidi.algorithm import get_display
    _HAS_PDF = True
except ImportError:
    _HAS_PDF = False

try:
    import fitz  # PyMuPDF
    _HAS_PYMUPDF = True
except ImportError:
    _HAS_PYMUPDF = False


# ══════════════════════════════════════════════════════════════════════════════
# TEXT EXTRACTION
# ══════════════════════════════════════════════════════════════════════════════

async def extract_full_text(file_bytes: bytes, filename: str, max_chars: int = 15_000) -> str:
    """Extract as much text as possible from the uploaded file."""
    fn = filename.lower()
    try:
        if fn.endswith(".pdf"):
            if _HAS_PYMUPDF:
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                pages = []
                for page in doc:
                    pages.append(page.get_text())
                return "\n".join(pages)[:max_chars]
            return ""

        if fn.endswith((".docx", ".doc")):
            if _HAS_DOCX:
                d = DocxDocument(io.BytesIO(file_bytes))
                return "\n".join(p.text for p in d.paragraphs)[:max_chars]
            return ""

        if fn.endswith((".xlsx", ".xls")):
            if _HAS_EXCEL:
                import pandas as pd
                dfs = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
                parts = []
                for sheet, df in dfs.items():
                    parts.append(f"=== Sheet: {sheet} ===")
                    parts.append(df.to_string(index=False, max_rows=200))
                return "\n\n".join(parts)[:max_chars]
            return ""

        if fn.endswith(".csv"):
            if _HAS_EXCEL:
                import pandas as pd
                df = pd.read_csv(io.BytesIO(file_bytes), nrows=500)
                return df.to_string(index=False)[:max_chars]
            return ""

        if fn.endswith((".pptx",)):
            if _HAS_PPTX:
                prs = Presentation(io.BytesIO(file_bytes))
                slides_text = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_parts = [f"--- Slide {i} ---"]
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_parts.append(shape.text.strip())
                    slides_text.append("\n".join(slide_parts))
                return "\n\n".join(slides_text)[:max_chars]
            return ""

        # Plain text fallback
        for enc in ("utf-8", "utf-8-sig", "windows-1256", "latin-1"):
            try:
                return file_bytes.decode(enc)[:max_chars]
            except UnicodeDecodeError:
                continue
        return file_bytes.decode("latin-1", errors="replace")[:max_chars]

    except Exception as exc:
        log.warning("smart_export: text extraction failed", filename=filename, error=str(exc))
        return f"[ملف: {filename}]"


# ══════════════════════════════════════════════════════════════════════════════
# AI ANALYSIS PROMPTS
# ══════════════════════════════════════════════════════════════════════════════

EXPORT_PROMPTS: Dict[str, Dict[str, str]] = {
    "excel": {
        "تحليل البيانات": "Extract all tabular data, statistics, and numbers. Return structured JSON with sheets, headers, and rows.",
        "تقرير مالي":     "Analyze financial data: revenues, costs, profits, KPIs. Structure for Excel with multiple sheets.",
        "ملخص تنفيذي":    "Extract key metrics and KPIs suitable for an executive Excel dashboard.",
    },
    "word": {
        "تقرير تنفيذي":   "Write a professional executive report in Arabic based on this document.",
        "تحليل شامل":     "Write a full detailed analysis report in Arabic with sections, findings, and recommendations.",
        "خطة عمل":        "Write an action plan in Arabic with objectives, tasks, timeline, and responsibilities.",
    },
    "pdf": {
        "تقرير رسمي":     "Generate a formal PDF report in Arabic with sections, analysis, and conclusions.",
        "ملخص تنفيذي":    "Generate a concise 1-2 page executive summary in Arabic.",
        "تقرير مفصل":     "Generate a comprehensive detailed report in Arabic.",
    },
    "powerpoint": {
        "عرض تقديمي":     "Create a 6-8 slide presentation in Arabic. Each slide: title + 3-5 bullet points.",
        "ملخص للإدارة":   "Create a 5-slide executive summary presentation in Arabic.",
        "تقرير القسم":    "Create a departmental report presentation with key findings.",
    },
    "powerbi": {
        "لوحة بيانات":    "Extract all numeric data and KPIs. Structure as Power BI compatible dataset.",
        "تحليل مالي":     "Extract financial metrics for Power BI financial dashboard.",
        "مؤشرات الأداء":  "Extract KPIs and performance indicators for Power BI.",
    },
}


def _build_analysis_prompt(text: str, filename: str, fmt: str, export_type: str) -> Tuple[str, str]:
    """Build system + user prompt for structured JSON response."""

    format_instructions = {
        "excel": textwrap.dedent("""
            Return ONLY valid JSON in this exact structure:
            {
              "title": "...",
              "summary": "...",
              "sheets": [
                {
                  "name": "sheet name",
                  "headers": ["col1", "col2", ...],
                  "rows": [["val1", "val2", ...], ...],
                  "chart": {"type": "bar"|"line"|null, "title": "..."}
                }
              ]
            }
            Ensure all rows have the same number of columns as headers.
            If data is textual (not tabular), create a summary sheet with Key|Value pairs.
        """),
        "word": textwrap.dedent("""
            Return ONLY valid JSON in this exact structure:
            {
              "title": "...",
              "subtitle": "...",
              "date": "...",
              "executive_summary": "...",
              "sections": [
                {
                  "heading": "...",
                  "paragraphs": ["paragraph 1", "paragraph 2"],
                  "bullet_points": ["point 1", "point 2"],
                  "table": {"headers": [...], "rows": [[...]]}
                }
              ],
              "conclusions": "...",
              "recommendations": ["rec 1", "rec 2"]
            }
            Write everything in Arabic. Be comprehensive and professional.
        """),
        "pdf": textwrap.dedent("""
            Return ONLY valid JSON in this exact structure:
            {
              "title": "...",
              "subtitle": "...",
              "date": "...",
              "executive_summary": "...",
              "sections": [
                {"heading": "...", "content": "...", "bullets": ["...", "..."]}
              ],
              "conclusions": "...",
              "recommendations": ["...", "..."]
            }
            Write everything in Arabic. Professional tone.
        """),
        "powerpoint": textwrap.dedent("""
            Return ONLY valid JSON in this exact structure:
            {
              "title": "...",
              "subtitle": "...",
              "slides": [
                {
                  "title": "slide title",
                  "layout": "title_content"|"two_column"|"bullets",
                  "bullets": ["point 1", "point 2", "point 3"],
                  "notes": "speaker notes"
                }
              ]
            }
            First slide = title slide. Last slide = conclusion/next steps.
            Write titles in Arabic. 6-8 slides total.
        """),
        "powerbi": textwrap.dedent("""
            Return ONLY valid JSON in this exact structure:
            {
              "dataset_name": "...",
              "description": "...",
              "tables": [
                {
                  "name": "TableName",
                  "columns": [{"name": "col", "dataType": "string"|"number"|"date"}],
                  "rows": [{"col": "value", ...}, ...]
                }
              ],
              "measures": [
                {"name": "Measure Name", "expression": "DAX expression", "description": "..."}
              ],
              "suggested_visuals": [
                {"type": "bar_chart"|"pie_chart"|"kpi_card"|"table", "title": "...", "fields": ["..."]}
              ]
            }
        """),
    }

    system = (
        f"You are an expert business analyst and data specialist. "
        f"Analyze the provided document and generate structured output for {fmt.upper()} export. "
        f"Export type requested: {export_type}. "
        f"CRITICAL: Return ONLY valid JSON. No markdown code blocks. No explanations outside the JSON. "
        f"All text fields should be in Arabic unless they are technical terms. "
    ) + format_instructions.get(fmt, "")

    user = (
        f"Document: {filename}\n"
        f"Export Format: {fmt.upper()}\n"
        f"Export Type: {export_type}\n\n"
        f"Document Content:\n{text}\n\n"
        f"Generate the structured {fmt.upper()} content now."
    )
    return system, user


async def ai_analyze(
    text: str,
    filename: str,
    fmt: str,
    export_type: str,
    llm,
) -> dict:
    """Call LLM to analyze document and return structured JSON for file generation."""
    system, user_prompt = _build_analysis_prompt(text, filename, fmt, export_type)

    response = await llm.generate(
        prompt=user_prompt,
        system=system,
        temperature=0.2,
        max_tokens=4096,
        trust_system=True,
    )

    raw = response.content.strip()
    # Strip markdown code fences if present
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)
    raw = raw.strip()

    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        log.warning("smart_export: JSON parse failed, attempting repair", error=str(e))
        # Try to extract JSON object/array
        match = re.search(r"\{.*\}", raw, re.DOTALL)
        if match:
            try:
                return json.loads(match.group())
            except Exception:
                pass
        # Return minimal fallback
        return {
            "title": f"تحليل: {filename}",
            "summary": raw[:2000],
            "sections": [{"heading": "المحتوى", "content": raw[:2000], "bullets": []}],
        }


# ══════════════════════════════════════════════════════════════════════════════
# GENERATORS
# ══════════════════════════════════════════════════════════════════════════════

# ─── Excel ─────────────────────────────────────────────────────────────────

EXCEL_HEADER_FILL = "1E3A5F"
EXCEL_ACCENT     = "3B82F6"
EXCEL_ALT_ROW    = "F0F4FF"

def generate_excel(data: dict) -> BytesIO:
    wb = openpyxl.Workbook()
    # Remove default sheet
    wb.remove(wb.active)

    sheets = data.get("sheets", [])
    if not sheets:
        # Build a simple summary sheet from whatever data exists
        sheets = [{
            "name": "ملخص",
            "headers": ["البند", "القيمة"],
            "rows": [[k, str(v)] for k, v in data.items() if not isinstance(v, (list, dict))],
        }]

    for sheet_data in sheets:
        ws = wb.create_sheet(title=str(sheet_data.get("name", "Sheet"))[:31])
        headers: List[str] = sheet_data.get("headers", [])
        rows: List[List[Any]] = sheet_data.get("rows", [])

        # Title row
        if data.get("title"):
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(len(headers), 1))
            title_cell = ws.cell(row=1, column=1, value=data["title"])
            title_cell.font = Font(bold=True, size=14, color="FFFFFF")
            title_cell.fill = PatternFill("solid", fgColor=EXCEL_HEADER_FILL)
            title_cell.alignment = Alignment(horizontal="center", vertical="center")
            ws.row_dimensions[1].height = 30
            header_start = 3
        else:
            header_start = 1

        # Headers
        if headers:
            for col_idx, h in enumerate(headers, 1):
                cell = ws.cell(row=header_start, column=col_idx, value=str(h))
                cell.font = Font(bold=True, size=11, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor=EXCEL_ACCENT)
                cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                ws.column_dimensions[cell.column_letter].width = max(15, len(str(h)) + 4)

            # Data rows
            for r_idx, row in enumerate(rows, header_start + 1):
                fill = PatternFill("solid", fgColor=EXCEL_ALT_ROW) if r_idx % 2 == 0 else None
                for c_idx, val in enumerate(row[:len(headers)], 1):
                    cell = ws.cell(row=r_idx, column=c_idx, value=val)
                    if fill:
                        cell.fill = fill
                    cell.alignment = Alignment(horizontal="right", vertical="center", wrap_text=True)
                    # Try numeric conversion
                    if isinstance(val, str):
                        try:
                            cell.value = float(val.replace(",", ""))
                            cell.number_format = "#,##0.00"
                        except (ValueError, AttributeError):
                            pass

            # Auto-filter
            if rows:
                ws.auto_filter.ref = ws.dimensions

        # Chart
        chart_cfg = sheet_data.get("chart", {})
        if chart_cfg and chart_cfg.get("type") and rows and len(headers) >= 2:
            from openpyxl.chart import BarChart, LineChart, Reference
            ChartCls = BarChart if chart_cfg["type"] == "bar" else BarChart
            chart = ChartCls()
            chart.title = chart_cfg.get("title", "")
            chart.style = 10
            chart.grouping = "clustered"
            chart.height = 12
            chart.width = 20

            data_ref = Reference(ws, min_col=2, max_col=len(headers),
                                 min_row=header_start, max_row=header_start + len(rows))
            cats = Reference(ws, min_col=1, min_row=header_start + 1, max_row=header_start + len(rows))
            chart.add_data(data_ref, titles_from_data=True)
            chart.set_categories(cats)
            # Place chart below data
            chart_row = header_start + len(rows) + 3
            ws.add_chart(chart, f"A{chart_row}")

    buf = BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


# ─── Word ──────────────────────────────────────────────────────────────────

def _docx_heading(doc: "DocxDocument", text: str, level: int = 1) -> None:
    h = doc.add_heading(text, level=level)
    h.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    for run in h.runs:
        run.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)


def _docx_para(doc: "DocxDocument", text: str, bold: bool = False) -> None:
    p = doc.add_paragraph(text)
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    if bold:
        for run in p.runs:
            run.bold = True


def generate_word(data: dict) -> BytesIO:
    doc = DocxDocument()

    # Document-wide RTL direction
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    # Title page
    title_para = doc.add_heading(data.get("title", "تقرير"), 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    if data.get("subtitle"):
        sub = doc.add_paragraph(data["subtitle"])
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER

    date_str = data.get("date") or datetime.now(timezone.utc).strftime("%Y-%m-%d")
    d = doc.add_paragraph(f"التاريخ: {date_str}")
    d.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()

    # Executive summary
    if data.get("executive_summary"):
        _docx_heading(doc, "الملخص التنفيذي", 1)
        _docx_para(doc, data["executive_summary"])
        doc.add_paragraph()

    # Sections
    for sec in data.get("sections", []):
        if sec.get("heading"):
            _docx_heading(doc, sec["heading"], 2)
        for para in sec.get("paragraphs", []):
            if para:
                _docx_para(doc, para)
        # Bullet points
        for bp in sec.get("bullet_points", []):
            if bp:
                p = doc.add_paragraph(style="List Bullet")
                p.add_run(bp)
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        # Inline table
        tbl = sec.get("table")
        if tbl and tbl.get("headers") and tbl.get("rows"):
            t = doc.add_table(rows=1, cols=len(tbl["headers"]))
            t.style = "Table Grid"
            hdr_cells = t.rows[0].cells
            for i, h in enumerate(tbl["headers"]):
                hdr_cells[i].text = str(h)
                hdr_cells[i].paragraphs[0].runs[0].bold = True
            for row in tbl["rows"]:
                cells = t.add_row().cells
                for i, v in enumerate(row[:len(tbl["headers"])]):
                    cells[i].text = str(v)
        doc.add_paragraph()

    # Conclusions
    if data.get("conclusions"):
        _docx_heading(doc, "الخلاصة والاستنتاجات", 1)
        _docx_para(doc, data["conclusions"])

    # Recommendations
    recs = data.get("recommendations", [])
    if recs:
        _docx_heading(doc, "التوصيات", 1)
        for rec in recs:
            p = doc.add_paragraph(style="List Number")
            p.add_run(rec)
            p.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# ─── PDF (ReportLab — Arabic-aware) ────────────────────────────────────────

def _reshape_ar(text: str) -> str:
    """Reshape Arabic text for correct RTL display in ReportLab."""
    try:
        reshaped = arabic_reshaper.reshape(str(text))
        return get_display(reshaped)
    except Exception:
        return str(text)


def generate_pdf(data: dict) -> BytesIO:
    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        rightMargin=2 * cm,
        leftMargin=2 * cm,
        topMargin=2.5 * cm,
        bottomMargin=2.5 * cm,
        title=data.get("title", "تقرير"),
    )

    styles = getSampleStyleSheet()
    # RTL paragraph style
    rtl_normal = ParagraphStyle(
        "RTL_Normal",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=11,
        leading=18,
        alignment=2,  # RIGHT
        spaceAfter=8,
    )
    rtl_h1 = ParagraphStyle(
        "RTL_H1",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=15,
        alignment=2,
        textColor=colors.HexColor("#1E3A5F"),
        spaceAfter=10,
        spaceBefore=14,
    )
    rtl_h2 = ParagraphStyle(
        "RTL_H2",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=12,
        alignment=2,
        textColor=colors.HexColor("#3B82F6"),
        spaceAfter=8,
        spaceBefore=10,
    )
    rtl_bullet = ParagraphStyle(
        "RTL_Bullet",
        parent=rtl_normal,
        leftIndent=20,
        bulletIndent=10,
    )

    story = []

    # Title
    title_text = _reshape_ar(data.get("title", "تقرير"))
    story.append(Paragraph(title_text, ParagraphStyle(
        "BigTitle", parent=rtl_h1, fontSize=20, alignment=1,  # CENTER
        textColor=colors.HexColor("#1E3A5F"), spaceAfter=6,
    )))

    if data.get("subtitle"):
        story.append(Paragraph(_reshape_ar(data["subtitle"]), ParagraphStyle(
            "Subtitle", parent=rtl_normal, alignment=1, textColor=colors.grey
        )))

    date_str = data.get("date") or datetime.now(timezone.utc).strftime("%Y-%m-%d")
    story.append(Paragraph(_reshape_ar(f"التاريخ: {date_str}"), ParagraphStyle(
        "DateStyle", parent=rtl_normal, alignment=1, textColor=colors.grey, fontSize=9
    )))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#3B82F6"),
                             spaceAfter=14, spaceBefore=6))

    # Executive summary
    if data.get("executive_summary"):
        story.append(Paragraph(_reshape_ar("الملخص التنفيذي"), rtl_h1))
        story.append(Paragraph(_reshape_ar(data["executive_summary"]), rtl_normal))
        story.append(Spacer(1, 10))

    # Sections
    for sec in data.get("sections", []):
        if sec.get("heading"):
            story.append(Paragraph(_reshape_ar(sec["heading"]), rtl_h2))
        if sec.get("content"):
            story.append(Paragraph(_reshape_ar(sec["content"]), rtl_normal))
        for b in sec.get("bullets", []):
            if b:
                story.append(Paragraph(f"• {_reshape_ar(b)}", rtl_bullet))
        story.append(Spacer(1, 8))

    # Conclusions
    if data.get("conclusions"):
        story.append(Paragraph(_reshape_ar("الخلاصة"), rtl_h1))
        story.append(Paragraph(_reshape_ar(data["conclusions"]), rtl_normal))

    # Recommendations
    recs = data.get("recommendations", [])
    if recs:
        story.append(Paragraph(_reshape_ar("التوصيات"), rtl_h1))
        for i, rec in enumerate(recs, 1):
            story.append(Paragraph(f"{i}. {_reshape_ar(rec)}", rtl_bullet))

    # Footer note
    story.append(Spacer(1, 20))
    story.append(HRFlowable(width="100%", thickness=0.5, color=colors.lightgrey))
    story.append(Paragraph(
        _reshape_ar("صدر بواسطة منصة ناطقة للذكاء الاصطناعي المؤسسي"),
        ParagraphStyle("Footer", parent=rtl_normal, fontSize=8, textColor=colors.grey, alignment=1)
    ))

    doc.build(story)
    buf.seek(0)
    return buf


# ─── PowerPoint ────────────────────────────────────────────────────────────

SLIDE_BG    = (6, 13, 26)     # #060d1a
SLIDE_TITLE = (59, 130, 246)  # #3b82f6
SLIDE_TEXT  = (204, 217, 239) # #ccd9ef
SLIDE_ACCENT= (16, 185, 129)  # #10b981

def _add_pptx_textbox(slide, text: str, left: float, top: float, width: float, height: float,
                      font_size: int = 20, bold: bool = False, color: tuple = SLIDE_TEXT,
                      align=PP_ALIGN.RIGHT) -> None:
    from pptx.util import Inches as PI, Pt as PPt
    txBox = slide.shapes.add_textbox(PI(left), PI(top), PI(width), PI(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PPt(font_size)
    run.font.bold = bold
    run.font.color.rgb = PptRGB(*color)


def generate_powerpoint(data: dict) -> BytesIO:
    prs = Presentation()
    prs.slide_width  = PI(13.33)
    prs.slide_height = PI(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    slides_data = data.get("slides", [])
    if not slides_data:
        slides_data = [{"title": data.get("title", "عرض"), "bullets": [data.get("summary", "")], "layout": "bullets"}]

    for i, sld in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)

        # Dark background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = PptRGB(*SLIDE_BG)

        # Slide number indicator
        _add_pptx_textbox(slide, f"{i+1}/{len(slides_data)}",
                          left=12.0, top=6.9, width=1.1, height=0.4,
                          font_size=9, color=(58, 84, 114), align=PP_ALIGN.RIGHT)

        # Accent line
        from pptx.util import Pt as PPt
        line = slide.shapes.add_shape(1, PI(0.5), PI(1.3), PI(12.3), PPt(2))  # shape type 1 = rectangle
        line.fill.solid()
        line.fill.fore_color.rgb = PptRGB(*SLIDE_TITLE)
        line.line.fill.background()

        if i == 0:
            # Title slide
            _add_pptx_textbox(slide, data.get("title", sld.get("title", "")),
                              left=1, top=2.5, width=11, height=1.2,
                              font_size=34, bold=True, color=SLIDE_TEXT, align=PP_ALIGN.CENTER)
            subtitle = data.get("subtitle") or sld.get("notes", "")
            if subtitle:
                _add_pptx_textbox(slide, subtitle,
                                  left=2, top=4.0, width=9, height=0.8,
                                  font_size=16, color=SLIDE_TITLE, align=PP_ALIGN.CENTER)
            _add_pptx_textbox(slide, "ناطقة | Natiqa",
                              left=0, top=6.7, width=13, height=0.5,
                              font_size=10, color=(42, 74, 110), align=PP_ALIGN.CENTER)
        else:
            # Content slide
            _add_pptx_textbox(slide, sld.get("title", ""),
                              left=0.5, top=0.4, width=12, height=0.85,
                              font_size=22, bold=True, color=SLIDE_TITLE, align=PP_ALIGN.RIGHT)

            bullets = sld.get("bullets", [])
            if bullets:
                # Build bullet text box
                txBox = slide.shapes.add_textbox(PI(0.5), PI(1.5), PI(12), PI(5.2))
                tf = txBox.text_frame
                tf.word_wrap = True
                for j, bullet in enumerate(bullets[:6]):
                    p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
                    p.alignment = PP_ALIGN.RIGHT
                    run = p.add_run()
                    run.text = f"• {bullet}"
                    run.font.size = PPt(17)
                    run.font.color.rgb = PptRGB(*SLIDE_TEXT)
                    p.space_after = PPt(8)

            # Speaker notes
            if sld.get("notes"):
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = sld["notes"]

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ─── Power BI JSON dataset ─────────────────────────────────────────────────

def generate_powerbi(data: dict) -> BytesIO:
    """
    Generate a Power BI importable JSON dataset.
    Format: Push Dataset API compatible JSON + metadata file.
    """
    output = {
        "dataset_name": data.get("dataset_name", "Natiqa Export"),
        "description": data.get("description", ""),
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "generated_by": "Natiqa AI Platform",
        "tables": data.get("tables", []),
        "measures": data.get("measures", []),
        "suggested_visuals": data.get("suggested_visuals", []),
        "schema_version": "1.0",
    }
    buf = BytesIO(json.dumps(output, ensure_ascii=False, indent=2).encode("utf-8"))
    buf.seek(0)
    return buf


# ══════════════════════════════════════════════════════════════════════════════
# MAIN ENTRY POINT
# ══════════════════════════════════════════════════════════════════════════════

MIME_TYPES = {
    "excel":       "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "word":        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pdf":         "application/pdf",
    "powerpoint":  "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "powerbi":     "application/json",
}

FILE_EXTENSIONS = {
    "excel":       "xlsx",
    "word":        "docx",
    "pdf":         "pdf",
    "powerpoint":  "pptx",
    "powerbi":     "json",
}


async def generate_smart_export(
    file_bytes: bytes,
    filename: str,
    output_format: str,
    export_type: str,
    llm,
) -> Tuple[BytesIO, str, str]:
    """
    Full pipeline: extract → analyze → generate.
    Returns: (buffer, mime_type, output_filename)
    """
    fmt = output_format.lower()
    if fmt not in MIME_TYPES:
        raise ValueError(f"Unsupported format: {fmt}. Choose: {list(MIME_TYPES.keys())}")

    log.info("smart_export: starting", filename=filename, format=fmt, type=export_type)

    # 1. Extract text
    text = await extract_full_text(file_bytes, filename)
    log.info("smart_export: extracted text", chars=len(text))

    # 2. AI analysis
    ai_result = await ai_analyze(text, filename, fmt, export_type, llm)
    log.info("smart_export: AI analysis complete", keys=list(ai_result.keys()))

    # 3. Generate file
    generators = {
        "excel":      generate_excel,
        "word":       generate_word,
        "pdf":        generate_pdf,
        "powerpoint": generate_powerpoint,
        "powerbi":    generate_powerbi,
    }

    buf = generators[fmt](ai_result)

    base = filename.rsplit(".", 1)[0]
    out_name = f"{base}_natiqa_export.{FILE_EXTENSIONS[fmt]}"
    mime = MIME_TYPES[fmt]

    log.info("smart_export: file generated", output=out_name, size=buf.getbuffer().nbytes)
    return buf, mime, out_name
