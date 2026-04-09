"""
╔══════════════════════════════════════════════════════════════════════╗
║  NATIQA — File Generation Engine                                     ║
╚══════════════════════════════════════════════════════════════════════╝
"""
import io
import os
import re
from typing import List, Optional
from datetime import datetime

import pandas as pd
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import arabic_reshaper
from bidi.algorithm import get_display

import structlog

log = structlog.get_logger()

# ─── Arabic Support ──────────────────────────────────────────────────
# Linux and Windows common paths
_FONT_PATHS = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    "C:/Windows/Fonts/arial.ttf",
]
ARABIC_FONT_PATH = None
for path in _FONT_PATHS:
    if os.path.exists(path):
        ARABIC_FONT_PATH = path
        break

if ARABIC_FONT_PATH:
    try:
        pdfmetrics.registerFont(TTFont('Arabic', ARABIC_FONT_PATH))
        log.info("Arabic font registered", path=ARABIC_FONT_PATH)
    except Exception as e:
        log.warning("Could not register Arabic font, using Helvetica", error=str(e))
else:
    log.warning("No TTF font found for Arabic PDF rendering. Falling back to Helvetica.")


def reshape_text(text: str) -> str:
    """Reshape and reorder Arabic text for PDF rendering."""
    if not any("\u0600" <= c <= "\u06FF" for c in text):
        return text
    reshaped_text = arabic_reshaper.reshape(text)
    return get_display(reshaped_text)


# ─── Content Parsing ─────────────────────────────────────────────────
def parse_tables(text: str) -> List[pd.DataFrame]:
    """Detect [TABLE] content and convert to DataFrames."""
    tables = []
    # Simple regex for Markdown tables or [TABLE] tags
    table_pattern = re.compile(r"\[TABLE\](.*?)(?=\[|$)", re.S)
    matches = table_pattern.findall(text)
    
    for m in matches:
        rows = [r.strip() for r in m.strip().split("\n") if "|" in r]
        if not rows: continue
        data = [[cell.strip() for cell in row.split("|")] for row in rows]
        # Clean empty cells from padding
        data = [[c for c in row if c] for row in data]
        if len(data) > 1:
            tables.append(pd.DataFrame(data[1:], columns=data[0]))
    return tables


def parse_slides(text: str) -> List[dict]:
    """Detect [SLIDE] tags."""
    slides = []
    slide_pattern = re.compile(r"\[SLIDE\](.*?)(?=\[|$)", re.S)
    matches = slide_pattern.findall(text)
    for m in matches:
        lines = [l.strip() for l in m.strip().split("\n") if l.strip()]
        if lines:
            slides.append({
                "title": lines[0],
                "content": lines[1:] if len(lines) > 1 else []
            })
    return slides


# ─── Generator Service ───────────────────────────────────────────────
class FileGenerator:
    
    @staticmethod
    def to_txt(content: str) -> bytes:
        return content.encode("utf-8")

    @staticmethod
    def to_pdf(content: str) -> bytes:
        buffer = io.BytesIO()
        c = canvas.Canvas(buffer, pagesize=A4)
        width, height = A4
        y = height - 50
        
        c.setFont("Arabic", 12)
        
        lines = content.split("\n")
        for line in lines:
            if y < 50:
                c.showPage()
                y = height - 50
                c.setFont("Arabic", 12)
            
            # Handle RTL alignment or simple centering for headings
            is_heading = "[HEADING]" in line
            text = line.replace("[HEADING]", "").strip()
            processed_text = reshape_text(text)
            
            if is_heading:
                c.setFont("Arabic", 16)
                c.drawCentredString(width/2, y, processed_text)
                y -= 30
                c.setFont("Arabic", 12)
            else:
                # Right align for Arabic
                c.drawRightString(width - 50, y, processed_text)
                y -= 20
                
        c.save()
        buffer.seek(0)
        return buffer.getvalue()

    @staticmethod
    def to_docx(content: str) -> bytes:
        doc = Document()
        # Better styling for Arabic
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Arial'
        font.size = Pt(12)
        
        lines = content.split("\n")
        for line in lines:
            if "[HEADING]" in line:
                doc.add_heading(line.replace("[HEADING]", "").strip(), level=1)
            elif "|" in line and "-|-" not in line: # Basic table detection
                # Ideally use the tables from parse_tables, but simple line-by-line is safer for mixed content
                doc.add_paragraph(line)
            else:
                doc.add_paragraph(line)
                
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    @staticmethod
    def to_xlsx(content: str) -> bytes:
        tables = parse_tables(content)
        buffer = io.BytesIO()
        with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
            if not tables:
                pd.DataFrame({"Content": [content]}).to_excel(writer, index=False, sheet_name="Report")
            for i, df in enumerate(tables):
                df.to_excel(writer, index=False, sheet_name=f"Table_{i+1}")
        buffer.seek(0)
        return buffer.getvalue()

    @staticmethod
    def to_csv(content: str) -> bytes:
        tables = parse_tables(content)
        if tables:
            return tables[0].to_csv(index=False).encode("utf-8")
        return pd.DataFrame({"Content": [content]}).to_csv(index=False).encode("utf-8")

    @staticmethod
    def to_pbi_csv(content: str) -> bytes:
        """
        Specialized Power BI Friendly CSV.
        - Clean headers (no special chars).
        - Standardized dates (YYYY-MM-DD).
        - UTF-8-BOM encoding.
        """
        tables = parse_tables(content)
        df = tables[0] if tables else pd.DataFrame({"Content": [content]})
        
        # 1. Clean Headers (alphanumeric + underscore only)
        df.columns = [re.sub(r'[^\w\s]', '', str(c)).strip().replace(' ', '_') for c in df.columns]
        
        # 2. Standardize Dates
        for col in df.columns:
            if any(k in col.lower() for k in ['date', 'time', 'تاريخ', 'وقت']):
                try:
                    df[col] = pd.to_datetime(df[col], errors='coerce').dt.strftime('%Y-%m-%d')
                except Exception:
                    pass
        
        # 3. Export with UTF-8-BOM for Excel/PowerBI compatibility with Arabic
        csv_data = df.to_csv(index=False).encode("utf-8")
        return b'\xef\xbb\xbf' + csv_data

    @staticmethod
    def to_pptx(content: str) -> bytes:
        prs = Presentation()
        slides_data = parse_slides(content)
        
        if not slides_data:
            # Create a simple default slide if no tags found
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "AI Report"
            slide.placeholders[1].text = content[:500]
        else:
            for s in slides_data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = s["title"]
                tf = slide.placeholders[1].text_frame
                for point in s["content"]:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
                    
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
