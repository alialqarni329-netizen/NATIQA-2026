from datetime import datetime, timezone
from io import BytesIO
from docx import Document as DocxDocument
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
import os

class ExportService:
    @staticmethod
    def generate_word_report(stats_data: dict, logo_path: str = None) -> BytesIO:
        doc = DocxDocument()
        
        # Header with Logo
        if logo_path and os.path.exists(logo_path):
            header = doc.sections[0].header
            header_para = header.paragraphs[0]
            run = header_para.add_run()
            run.add_picture(logo_path, width=Inches(1.5))
            header_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

        # Title
        title = doc.add_heading('تقرير أداء منصة ناطقة (Natiqa)', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Date
        date_para = doc.add_paragraph()
        date_para.add_run(f"تاريخ التقرير: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M:%S')} (UTC)")
        date_para.alignment = WD_ALIGN_PARAGRAPH.LEFT

        doc.add_heading('ملخص مؤشرات الأداء الحالية', level=1)
        
        # Stats Table
        table = doc.add_table(rows=1, cols=2)
        table.style = 'Table Grid'
        header_cells = table.rows[0].cells
        header_cells[0].text = 'المؤشر'
        header_cells[1].text = 'القيمة'
        
        stats_map = {
            "إجمالي المستخدمين": stats_data.get("users", {}).get("total", 0),
            "المستخدمين النشطين": stats_data.get("users", {}).get("active", 0),
            "إجمالي المؤسسات": stats_data.get("organizations", {}).get("total", 0),
            "المؤسسات النشطة": stats_data.get("organizations", {}).get("active", 0),
            "إجمالي المستندات المعالجة": stats_data.get("documents", {}).get("total", 0),
        }
        
        for key, value in stats_map.items():
            row_cells = table.add_row().cells
            row_cells[0].text = key
            row_cells[1].text = str(value)

        doc.add_paragraph("\nتم تقديم هذا التقرير آلياً بواسطة نظام ناطقة لإدارة الذكاء الاصطناعي المؤسسي.")
        
        target = BytesIO()
        doc.save(target)
        target.seek(0)
        return target

    @staticmethod
    def generate_pptx_presentation(stats_data: dict, logo_path: str = None) -> BytesIO:
        prs = Presentation()
        
        # Slide 1: Title
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "منصة ناطقة NATIQA"
        subtitle.text = "تقرير حالة المنصة ومؤشرات النمو\nEnterprise AI Platform Status"

        # Slide 2: Key Metrics
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "المؤشرات الرئيسية (Key KPIs)"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = f"• إجمالي المستخدمين: {stats_data.get('users', {}).get('total', 0)}"
        tf.add_paragraph().text = f"• إجمالي الشركات المسجلة: {stats_data.get('organizations', {}).get('total', 0)}"
        tf.add_paragraph().text = f"• المستندات المعالجة: {stats_data.get('documents', {}).get('total', 0)}"
        
        # Add logo to every slide if exists
        if logo_path and os.path.exists(logo_path):
            for s in prs.slides:
                s.shapes.add_picture(logo_path, PptInches(0.2), PptInches(0.2), width=PptInches(1))

        # Slide 3: Growth Insight
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "نظرة على نمو المنصة"
        
        content = slide.placeholders[1]
        content.text = "تشهد المنصة نمواً مستمراً في قاعدة البيانات وتحليل المستندات، مع زيادة مطردة في طلبات المؤسسات على الذكاء الاصطناعي السيادي والآمن."

        target = BytesIO()
        prs.save(target)
        target.seek(0)
        return target
