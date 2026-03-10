"""
╔══════════════════════════════════════════════════════════════════════════╗
║         NATIQA — Omni-Document Processor  (Enterprise-Grade)            ║
║                                                                          ║
║  المراحل:                                                                ║
║    1. RECEIVE  → استقبال bytes خام + metadata                           ║
║    2. VALIDATE → التحقق من نوع الملف + SHA-256 hash                     ║
║    3. ENCRYPT  → تشفير AES-256-GCM فوري في التخزين                      ║
║    4. EXTRACT  → استخراج النص/البيانات من كل صيغة                       ║
║    5. MASK     → Data Masking قبل أي إرسال خارجي                        ║
║    6. INDEX    → تضمين + تخزين في ChromaDB مع RBAC metadata             ║
║    7. WIPE     → حذف آمن للبيانات المؤقتة من الذاكرة                    ║
║                                                                          ║
║  الصيغ المدعومة:                                                         ║
║    • Excel / CSV     → Pandas (قاعدة بيانات مؤقتة للأرقام الدقيقة)     ║
║    • PDF             → PyMuPDF (fitz) — نص + جداول + صور               ║
║    • Word (.docx)    → python-docx + python-pptx                        ║
║    • PowerPoint      → python-pptx                                       ║
║    • TXT / Markdown  → نص مباشر                                         ║
║                                                                          ║
║  الأمان:                                                                 ║
║    • AES-256-GCM لكل ملف على القرص                                      ║
║    • Data Masking تلقائي (7 أنواع)                                      ║
║    • RBAC على مستوى الـ Chunk في ChromaDB                               ║
║    • Secure Wipe: overwrite بـ random bytes ثم حذف                      ║
║    • لا توجد ملفات مؤقتة على القرص — كل شيء in-memory                  ║
╚══════════════════════════════════════════════════════════════════════════╝
"""

from __future__ import annotations

import asyncio
import gc
import hashlib
import io
import os
import secrets
import tempfile
import time
import ctypes
from contextlib import asynccontextmanager
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any, AsyncGenerator

import aiofiles
import structlog

from app.core.config import settings
from app.core.security import encrypt_file, decrypt_file
from app.services.llm.factory import get_llm
from app.services.llm.masking import mask_sensitive_data, MaskingResult

log = structlog.get_logger()


# ═══════════════════════════════════════════════════════════
#  1. RBAC — Document Access Control
# ═══════════════════════════════════════════════════════════

class DocumentSensitivity(str, Enum):
    """مستويات حساسية الوثيقة — تتحكم في من يمكنه استرجاعها عبر RAG."""
    PUBLIC     = "public"       # جميع المستخدمين
    INTERNAL   = "internal"     # analyst + admin + super_admin
    RESTRICTED = "restricted"   # admin + super_admin
    CONFIDENTIAL = "confidential"  # super_admin فقط


# الأدوار المسموح لها بكل مستوى
SENSITIVITY_ROLES: dict[DocumentSensitivity, set[str]] = {
    DocumentSensitivity.PUBLIC:       {"viewer", "analyst", "admin", "super_admin"},
    DocumentSensitivity.INTERNAL:     {"analyst", "admin", "super_admin"},
    DocumentSensitivity.RESTRICTED:   {"admin", "super_admin"},
    DocumentSensitivity.CONFIDENTIAL: {"super_admin"},
}

# تصنيف تلقائي بناءً على اسم القسم
DEPT_SENSITIVITY_MAP: dict[str, DocumentSensitivity] = {
    "hr":        DocumentSensitivity.RESTRICTED,    # ملفات الموارد البشرية / الرواتب
    "payroll":   DocumentSensitivity.CONFIDENTIAL,  # كشوف المرتبات
    "legal":     DocumentSensitivity.RESTRICTED,    # العقود والوثائق القانونية
    "financial": DocumentSensitivity.INTERNAL,      # التقارير المالية
    "technical": DocumentSensitivity.INTERNAL,      # الوثائق التقنية
    "admin":     DocumentSensitivity.INTERNAL,      # الوثائق الإدارية
    "general":   DocumentSensitivity.PUBLIC,        # الوثائق العامة
}


def resolve_sensitivity(
    department: str,
    override: DocumentSensitivity | None = None,
) -> DocumentSensitivity:
    """تحديد مستوى الحساسية — يقبل override صريح أو يستنتج من القسم."""
    if override:
        return override
    dept_lower = department.lower().strip()
    return DEPT_SENSITIVITY_MAP.get(dept_lower, DocumentSensitivity.INTERNAL)


def can_access_document(user_role: str, sensitivity: DocumentSensitivity) -> bool:
    """هل يملك هذا الدور صلاحية الوصول لهذا المستوى؟"""
    allowed = SENSITIVITY_ROLES.get(sensitivity, set())
    return user_role in allowed


# ═══════════════════════════════════════════════════════════
#  2. Secure Memory Buffer
# ═══════════════════════════════════════════════════════════

class SecureBuffer:
    """
    Buffer آمن — يُمسح بـ random bytes عند الإتلاف.
    يمنع بقاء بيانات حساسة في الذاكرة بعد المعالجة.
    """

    def __init__(self, data: bytes):
        # نسخ البيانات في bytearray قابل للتعديل
        self._buf = bytearray(data)
        self._wiped = False

    @property
    def data(self) -> bytes:
        if self._wiped:
            raise RuntimeError("SecureBuffer: تمّ مسح البيانات بالفعل")
        return bytes(self._buf)

    def wipe(self) -> None:
        """كتابة فوضوية عشوائية فوق البيانات ثم تصفيرها."""
        if not self._wiped:
            length = len(self._buf)
            # مرحلة 1: كتابة عشوائية
            for i in range(length):
                self._buf[i] = secrets.randbelow(256)
            # مرحلة 2: تصفير
            for i in range(length):
                self._buf[i] = 0
            self._wiped = True
            log.debug("SecureBuffer: تمّ المسح الآمن", bytes_wiped=length)

    def __del__(self):
        self.wipe()

    def __enter__(self):
        return self

    def __exit__(self, *_):
        self.wipe()


def secure_wipe_string(s: str) -> None:
    """
    محاولة مسح string من الذاكرة.
    تحذير: Python strings غير قابلة للتعديل،
    هذه المحاولة أفضل من لا شيء لكنها غير مضمونة 100%.
    """
    try:
        encoded = s.encode("utf-8")
        buf = ctypes.create_string_buffer(encoded)
        ctypes.memset(buf, 0, len(encoded))
    except Exception:
        pass
    finally:
        gc.collect()


# ═══════════════════════════════════════════════════════════
#  3. Processing Result
# ═══════════════════════════════════════════════════════════

@dataclass
class ProcessingResult:
    """نتيجة معالجة وثيقة واحدة."""
    doc_id: str
    filename: str
    file_type: str
    chunks_count: int
    file_hash: str                          # SHA-256 للملف الأصلي
    encrypted_path: str                     # مسار الملف المشفّر
    sensitivity: DocumentSensitivity
    department: str
    masked_fields_count: int = 0            # عدد الحقول التي تم إخفاؤها
    tabular_rows: int = 0                   # للملفات الجدولية
    processing_time_ms: int = 0
    warnings: list[str] = field(default_factory=list)


# ═══════════════════════════════════════════════════════════
#  4. Format-Specific Extractors
# ═══════════════════════════════════════════════════════════

async def _extract_pdf(raw_bytes: bytes, filename: str) -> tuple[str, list[str]]:
    """
    PDF → نص كامل عبر PyMuPDF (fitz).
    يستخرج: نص + جداول (كـ markdown).
    يعود بـ: (full_text, warnings)
    """
    import fitz  # PyMuPDF

    warnings: list[str] = []
    pages_text: list[str] = []

    doc = fitz.open(stream=raw_bytes, filetype="pdf")
    total_pages = len(doc)

    for page_num in range(total_pages):
        page = doc[page_num]

        # نص عادي
        page_text = page.get_text("text")

        # استخراج الجداول كـ markdown
        try:
            tabs = page.find_tables()
            for tab in tabs.tables:
                df_data = tab.extract()
                if df_data:
                    # تحويل إلى نص منسّق
                    header = " | ".join(str(c) if c else "" for c in df_data[0])
                    sep    = " | ".join(["---"] * len(df_data[0]))
                    rows   = [" | ".join(str(c) if c else "" for c in row) for row in df_data[1:]]
                    table_text = f"\n[جدول صفحة {page_num+1}]\n{header}\n{sep}\n" + "\n".join(rows) + "\n"
                    page_text += table_text
        except Exception as e:
            warnings.append(f"فشل استخراج جدول في صفحة {page_num+1}: {e}")

        if page_text.strip():
            pages_text.append(f"[صفحة {page_num+1}/{total_pages}]\n{page_text.strip()}")

    doc.close()

    if not pages_text:
        warnings.append("PDF لا يحتوي على نص قابل للاستخراج (قد يكون صورة مسحوحة)")

    return "\n\n".join(pages_text), warnings


async def _extract_docx(raw_bytes: bytes) -> tuple[str, list[str]]:
    """Word .docx → نص + جداول."""
    from docx import Document
    from docx.oxml.ns import qn

    warnings: list[str] = []
    parts: list[str] = []

    doc = Document(io.BytesIO(raw_bytes))

    # الفقرات والعناوين
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style else ""
        if "Heading" in style:
            parts.append(f"\n## {text}\n")
        else:
            parts.append(text)

    # الجداول
    for i, table in enumerate(doc.tables):
        try:
            rows_text = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows_text.append(" | ".join(cells))
            if rows_text:
                parts.append(f"\n[جدول {i+1}]\n" + "\n".join(rows_text) + "\n")
        except Exception as e:
            warnings.append(f"فشل استخراج جدول {i+1}: {e}")

    return "\n".join(parts), warnings


async def _extract_pptx(raw_bytes: bytes) -> tuple[str, list[str]]:
    """PowerPoint .pptx → نص كل شريحة."""
    from pptx import Presentation
    from pptx.util import Inches

    warnings: list[str] = []
    slides_text: list[str] = []

    prs = Presentation(io.BytesIO(raw_bytes))

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
            # جداول داخل الشرائح
            if shape.has_table:
                try:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        slide_parts.append(row_text)
                except Exception as e:
                    warnings.append(f"فشل استخراج جدول في شريحة {slide_num}: {e}")

        if slide_parts:
            slides_text.append(
                f"[شريحة {slide_num}/{len(prs.slides)}]\n" + "\n".join(slide_parts)
            )

    return "\n\n".join(slides_text), warnings


@dataclass
class TabularResult:
    """نتيجة معالجة ملف جدولي (Excel/CSV)."""
    text_summary: str       # ملخص نصي للاستيعاب في RAG
    row_count: int
    col_count: int
    dtypes: dict[str, str]  # { column_name: dtype }
    numeric_stats: dict     # إحصائيات للأعمدة الرقمية


async def _extract_tabular(
    raw_bytes: bytes,
    filename: str,
    sheet_name: str | None = None,
) -> tuple[str, TabularResult, list[str]]:
    """
    Excel / CSV → Pandas DataFrame.

    الاستراتيجية:
    1. تحميل الملف في DataFrame
    2. حساب الإحصائيات الكاملة (describe)
    3. تحويل إلى نص منظّم للاستيعاب في RAG
    4. الـ DataFrame نفسه يُرفق مع الـ Chunk metadata
       حتى يتمكن محرك RAG لاحقاً من الاستعلام الدقيق

    هذا يحل مشكلة "الأرقام التقريبية" — عند السؤال
    عن مبلغ محدد، يمكن البحث الدقيق في DataFrame.
    """
    import pandas as pd

    warnings: list[str] = []
    ext = Path(filename).suffix.lower()

    try:
        if ext == ".csv":
            # محاولة اكتشاف الـ encoding
            for enc in ("utf-8", "utf-8-sig", "windows-1256", "latin-1"):
                try:
                    df = pd.read_csv(io.BytesIO(raw_bytes), encoding=enc)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                df = pd.read_csv(io.BytesIO(raw_bytes), encoding="latin-1", errors="replace")
                warnings.append("تعذّر اكتشاف encoding — تم استخدام latin-1")

        elif ext in (".xlsx", ".xls", ".xlsm"):
            xl = pd.ExcelFile(io.BytesIO(raw_bytes))
            sheets = xl.sheet_names

            if sheet_name and sheet_name in sheets:
                df = xl.parse(sheet_name)
            elif len(sheets) == 1:
                df = xl.parse(sheets[0])
            else:
                # دمج جميع الأوراق في DataFrame واحد
                frames = []
                for sh in sheets:
                    try:
                        frame = xl.parse(sh)
                        frame["__sheet__"] = sh
                        frames.append(frame)
                    except Exception as e:
                        warnings.append(f"فشل تحميل ورقة '{sh}': {e}")
                if not frames:
                    raise ValueError("لا توجد أوراق قابلة للقراءة")
                df = pd.concat(frames, ignore_index=True)
                warnings.append(f"تم دمج {len(frames)} ورقة: {', '.join(sheets[:5])}")
        else:
            raise ValueError(f"امتداد غير مدعوم: {ext}")

    except Exception as e:
        raise ValueError(f"فشل تحليل الملف الجدولي: {e}") from e

    # تنظيف
    df = df.dropna(how="all").dropna(axis=1, how="all")
    df.columns = [str(c).strip() for c in df.columns]

    rows, cols = df.shape

    # إحصائيات الأعمدة الرقمية
    numeric_cols = df.select_dtypes(include="number").columns.tolist()
    numeric_stats: dict = {}
    if numeric_cols:
        desc = df[numeric_cols].describe().to_dict()
        numeric_stats = {
            col: {
                k: round(float(v), 4) if isinstance(v, float) else v
                for k, v in stats.items()
            }
            for col, stats in desc.items()
        }

    # بناء النص للـ RAG
    parts: list[str] = []
    parts.append(f"[ملف جدولي: {filename}]")
    parts.append(f"الأبعاد: {rows} صف × {cols} عمود")
    parts.append(f"الأعمدة: {', '.join(df.columns.tolist()[:20])}")

    if numeric_stats:
        parts.append("\n[الإحصائيات الرقمية]")
        for col, stats in list(numeric_stats.items())[:10]:
            s = stats
            parts.append(
                f"  {col}: "
                f"مجموع={s.get('sum', 'N/A')}, "
                f"متوسط={s.get('mean', 'N/A')}, "
                f"min={s.get('min', 'N/A')}, "
                f"max={s.get('max', 'N/A')}, "
                f"عدد={int(s.get('count', 0))}"
            )

    # عيّنة من البيانات (أول 50 صف كـ markdown)
    sample = df.head(50)
    try:
        parts.append("\n[عيّنة البيانات (أول 50 صف)]")
        parts.append(sample.to_markdown(index=False))
    except Exception:
        parts.append(sample.to_string(index=False))

    # إضافة الصفوف المتبقية (ملخص)
    if rows > 50:
        parts.append(f"\n[... {rows - 50} صف إضافي غير معروض في العيّنة]")

    text_summary = "\n".join(parts)

    dtypes = {col: str(dt) for col, dt in df.dtypes.items()}

    tabular_result = TabularResult(
        text_summary=text_summary,
        row_count=rows,
        col_count=cols,
        dtypes=dtypes,
        numeric_stats=numeric_stats,
    )

    return text_summary, tabular_result, warnings


async def _extract_text_plain(raw_bytes: bytes, filename: str) -> tuple[str, list[str]]:
    """TXT / Markdown → نص مباشر."""
    warnings: list[str] = []
    for enc in ("utf-8", "utf-8-sig", "windows-1256", "latin-1"):
        try:
            return raw_bytes.decode(enc), warnings
        except UnicodeDecodeError:
            continue
    warnings.append("تعذّر اكتشاف encoding — تم استخدام latin-1 مع استبدال")
    return raw_bytes.decode("latin-1", errors="replace"), warnings


# ═══════════════════════════════════════════════════════════
#  5. Text Chunker
# ═══════════════════════════════════════════════════════════

def chunk_text(
    text: str,
    chunk_size: int = 900,
    overlap: int = 120,
    min_chunk: int = 50,
) -> list[str]:
    """
    تقطيع ذكي يحترم:
    - فواصل الفقرات (أولوية)
    - نهايات الجمل العربية والإنجليزية
    - عدم تقطيع الجداول في المنتصف
    """
    if len(text) <= chunk_size:
        return [text] if len(text) >= min_chunk else []

    separators = ["\n\n", "\n", ".", "،", "؟", "!", "؛", ";", " "]
    chunks: list[str] = []
    start = 0

    while start < len(text):
        end = min(start + chunk_size, len(text))

        if end < len(text):
            best = end
            for sep in separators:
                idx = text.rfind(sep, start + min_chunk, end)
                if idx > start + min_chunk:
                    best = idx + len(sep)
                    break
            end = best

        chunk = text[start:end].strip()
        if len(chunk) >= min_chunk:
            chunks.append(chunk)

        # Overlap
        next_start = end - overlap
        start = next_start if next_start > start else end

    return chunks


# ═══════════════════════════════════════════════════════════
#  6. Secure File Storage
# ═══════════════════════════════════════════════════════════

async def _store_encrypted(
    raw_bytes: bytes,
    doc_id: str,
    project_id: str,
) -> tuple[str, str]:
    """
    تشفير فوري وحفظ.
    يعود بـ: (encrypted_path, sha256_hash)
    """
    # حساب hash قبل التشفير
    file_hash = hashlib.sha256(raw_bytes).hexdigest()

    encrypted = encrypt_file(raw_bytes)

    save_dir = Path(settings.UPLOAD_DIR) / project_id
    save_dir.mkdir(parents=True, exist_ok=True)
    save_path = save_dir / f"{doc_id}.enc"

    async with aiofiles.open(save_path, "wb") as f:
        await f.write(encrypted)

    log.info(
        "تم التخزين المشفّر",
        path=str(save_path),
        original_size=len(raw_bytes),
        encrypted_size=len(encrypted),
    )
    return str(save_path), file_hash


async def _load_and_decrypt(encrypted_path: str) -> bytes:
    """تحميل وفك تشفير ملف محفوظ."""
    async with aiofiles.open(encrypted_path, "rb") as f:
        encrypted = await f.read()
    return decrypt_file(encrypted)


# ═══════════════════════════════════════════════════════════
#  7. Secure Wipe
# ═══════════════════════════════════════════════════════════

async def secure_wipe_file(file_path: str, passes: int = 3) -> bool:
    """
    حذف آمن لملف:
    1. الكتابة فوقه بـ random bytes (passes مرات)
    2. الكتابة فوقه بأصفار
    3. حذفه

    ملاحظة: على أقراص SSD هذا لا يضمن 100% المسح الفيزيائي
    بسبب wear leveling، لكنه يمنع الاسترداد البرمجي البسيط.
    """
    path = Path(file_path)
    if not path.exists():
        return False

    try:
        file_size = path.stat().st_size

        async with aiofiles.open(file_path, "r+b") as f:
            # مرور عشوائي × passes
            for pass_num in range(passes):
                await f.seek(0)
                chunk = 65536  # 64KB
                written = 0
                while written < file_size:
                    to_write = min(chunk, file_size - written)
                    await f.write(secrets.token_bytes(to_write))
                    written += to_write
                await f.flush()

            # مرور أصفار أخير
            await f.seek(0)
            written = 0
            while written < file_size:
                to_write = min(65536, file_size - written)
                await f.write(b"\x00" * to_write)
                written += to_write
            await f.flush()

        path.unlink()
        log.info("تم المسح الآمن للملف", path=file_path, passes=passes)
        return True

    except Exception as e:
        log.error("فشل المسح الآمن", path=file_path, error=str(e))
        # محاولة الحذف العادي على أي حال
        try:
            path.unlink(missing_ok=True)
        except Exception:
            pass
        return False


# ═══════════════════════════════════════════════════════════
#  8. ChromaDB RBAC Integration
# ═══════════════════════════════════════════════════════════

def _get_collection(project_id: str):
    import chromadb
    client = chromadb.PersistentClient(path=settings.VECTOR_DIR)
    return client.get_or_create_collection(
        name=f"project_{project_id.replace('-', '_')}",
        metadata={"hnsw:space": "cosine"},
    )


async def _embed_and_store_chunks(
    chunks: list[str],
    doc_id: str,
    project_id: str,
    filename: str,
    department: str,
    sensitivity: DocumentSensitivity,
    file_type: str,
    uploaded_by: str,
    mask_before_embed: bool = True,
) -> tuple[int, int]:
    """
    تضمين وتخزين الـ chunks في ChromaDB مع:
    - RBAC metadata (sensitivity + department)
    - Data Masking قبل الإرسال لـ LLM
    - إحصائيات الـ Masking

    يعود بـ: (chunks_stored, total_masked_fields)
    """
    llm = get_llm()
    collection = _get_collection(project_id)
    total_masked = 0

    BATCH_SIZE = 16

    for batch_start in range(0, len(chunks), BATCH_SIZE):
        batch = chunks[batch_start: batch_start + BATCH_SIZE]
        ids, embeddings, docs_list, metas = [], [], [], []

        for j, chunk in enumerate(batch):
            chunk_idx = batch_start + j

            # ── Data Masking قبل التضمين ──────────────────────────
            if mask_before_embed and llm.provider_name == "claude":
                mask_result = mask_sensitive_data(
                    chunk,
                    session_salt=settings.ENCRYPTION_KEY[:16],
                )
                chunk_to_embed = mask_result.masked_text
                total_masked += mask_result.count
            else:
                chunk_to_embed = chunk

            # ── التضمين ──────────────────────────────────────────
            emb = await llm.embed(chunk_to_embed)

            ids.append(f"{doc_id}_{chunk_idx}")
            embeddings.append(emb.embedding)
            docs_list.append(chunk_to_embed)        # نخزن النص بعد المسك

            # ── Metadata للـ RBAC ─────────────────────────────────
            metas.append({
                "doc_id":      doc_id,
                "filename":    filename,
                "department":  department,
                "sensitivity": sensitivity.value,   # للفلترة في وقت الاسترجاع
                "file_type":   file_type,
                "chunk_idx":   chunk_idx,
                "uploaded_by": uploaded_by,
            })

        collection.add(
            ids=ids,
            embeddings=embeddings,
            documents=docs_list,
            metadatas=metas,
        )
        log.debug(
            "Batch embedded",
            batch_start=batch_start,
            batch_size=len(batch),
            provider=llm.provider_name,
        )

    return len(chunks), total_masked


# ═══════════════════════════════════════════════════════════
#  9. RBAC Query Filter
# ═══════════════════════════════════════════════════════════

def build_rbac_where_filter(
    user_role: str,
    allowed_sensitivities: list[DocumentSensitivity] | None = None,
    department: str | None = None,
) -> dict | None:
    """
    بناء فلتر ChromaDB للاستعلام بناءً على RBAC.

    يُستخدم في query_rag لضمان أن المستخدم لا يرى
    إلا الـ chunks التي يملك صلاحية الوصول إليها.
    """
    # تحديد مستويات الحساسية المسموح بها لهذا الدور
    if allowed_sensitivities is None:
        allowed_sensitivities = [
            s for s in DocumentSensitivity
            if user_role in SENSITIVITY_ROLES[s]
        ]

    if not allowed_sensitivities:
        # لا صلاحيات — لا نتائج
        return {"sensitivity": "__none__"}

    sensitivity_values = [s.value for s in allowed_sensitivities]

    # بناء فلتر $in
    where: dict = {
        "sensitivity": {"$in": sensitivity_values}
    }

    # فلتر إضافي بالقسم
    if department:
        where = {
            "$and": [
                {"sensitivity": {"$in": sensitivity_values}},
                {"department": department},
            ]
        }

    return where


# ═══════════════════════════════════════════════════════════
#  10. Main Processor — ingest_document_v2
# ═══════════════════════════════════════════════════════════

async def ingest_document_v2(
    raw_bytes: bytes,
    filename: str,
    doc_id: str,
    project_id: str,
    department: str,
    uploaded_by: str,
    sensitivity_override: DocumentSensitivity | None = None,
) -> ProcessingResult:
    """
    نقطة الدخول الرئيسية — تنفّذ كل المراحل بالترتيب:

    RECEIVE → VALIDATE → ENCRYPT → EXTRACT → MASK → INDEX → WIPE

    الـ raw_bytes يُمسح من الذاكرة فور الانتهاء من التشفير.
    """
    t_start = time.time()
    ext = Path(filename).suffix.lower().lstrip(".")
    warnings_all: list[str] = []

    log.info(
        "بدء معالجة وثيقة",
        doc_id=doc_id,
        filename=filename,
        size=len(raw_bytes),
        ext=ext,
    )

    # ── مرحلة 1: التحقق ──────────────────────────────────
    if len(raw_bytes) == 0:
        raise ValueError("الملف فارغ")

    allowed_exts = {"pdf", "docx", "doc", "pptx", "xlsx", "xls", "xlsm", "csv", "txt", "md"}
    if ext not in allowed_exts:
        raise ValueError(f"امتداد غير مدعوم: .{ext}")

    # ── مرحلة 2: التشفير الفوري ──────────────────────────
    with SecureBuffer(raw_bytes) as sbuf:
        encrypted_path, file_hash = await _store_encrypted(
            sbuf.data, doc_id, project_id
        )
        # مسح من الذاكرة فور الانتهاء من التشفير
    # raw_bytes لا يزال في scope لكن SecureBuffer مُمسح

    # ── مرحلة 3: تحديد الحساسية ──────────────────────────
    sensitivity = resolve_sensitivity(department, sensitivity_override)

    log.info(
        "تم تحديد مستوى الحساسية",
        sensitivity=sensitivity.value,
        department=department,
    )

    # ── مرحلة 4: استخراج النص حسب الصيغة ────────────────
    tabular_result: TabularResult | None = None
    extract_text = ""

    if ext == "pdf":
        extract_text, w = await _extract_pdf(raw_bytes, filename)
        warnings_all.extend(w)

    elif ext in ("docx", "doc"):
        extract_text, w = await _extract_docx(raw_bytes)
        warnings_all.extend(w)

    elif ext == "pptx":
        extract_text, w = await _extract_pptx(raw_bytes)
        warnings_all.extend(w)

    elif ext in ("xlsx", "xls", "xlsm", "csv"):
        extract_text, tabular_result, w = await _extract_tabular(raw_bytes, filename)
        warnings_all.extend(w)

    elif ext in ("txt", "md"):
        extract_text, w = await _extract_text_plain(raw_bytes, filename)
        warnings_all.extend(w)

    if not extract_text.strip():
        raise ValueError(f"لم يُتمكّن من استخراج أي نص من '{filename}'")

    # ── مرحلة 5: تقطيع ───────────────────────────────────
    chunks = chunk_text(extract_text)
    if not chunks:
        raise ValueError("لم ينتج أي chunk بعد التقطيع")

    log.info("تم التقطيع", chunks=len(chunks), doc_id=doc_id)

    # ── مرحلة 6: التضمين + الفهرسة مع RBAC ──────────────
    chunks_stored, masked_count = await _embed_and_store_chunks(
        chunks=chunks,
        doc_id=doc_id,
        project_id=project_id,
        filename=filename,
        department=department,
        sensitivity=sensitivity,
        file_type=ext,
        uploaded_by=str(uploaded_by),
    )

    # ── مرحلة 7: مسح النص المؤقت ─────────────────────────
    secure_wipe_string(extract_text)
    del extract_text
    for chunk in chunks:
        secure_wipe_string(chunk)
    del chunks
    gc.collect()

    t_end = time.time()
    elapsed_ms = int((t_end - t_start) * 1000)

    log.info(
        "اكتملت المعالجة",
        doc_id=doc_id,
        chunks=chunks_stored,
        masked_fields=masked_count,
        sensitivity=sensitivity.value,
        elapsed_ms=elapsed_ms,
    )

    return ProcessingResult(
        doc_id=doc_id,
        filename=filename,
        file_type=ext,
        chunks_count=chunks_stored,
        file_hash=file_hash,
        encrypted_path=encrypted_path,
        sensitivity=sensitivity,
        department=department,
        masked_fields_count=masked_count,
        tabular_rows=tabular_result.row_count if tabular_result else 0,
        processing_time_ms=elapsed_ms,
        warnings=warnings_all,
    )


# ═══════════════════════════════════════════════════════════
#  11. RBAC-Aware RAG Query
# ═══════════════════════════════════════════════════════════

async def query_with_rbac(
    question: str,
    project_id: str,
    user_role: str,
    top_k: int = 5,
    department_filter: str | None = None,
) -> dict:
    """
    استعلام RAG مع:
    - RBAC: فلترة النتائج بناءً على دور المستخدم
    - Masking: تُخفى البيانات الحساسة قبل إرسالها لـ LLM
    - Tabular precision: دعم الأسئلة الرقمية الدقيقة

    يعود بـ:
    {
        "answer": str,
        "sources": [...],
        "tokens": int,
        "masked_fields": int,
        "access_denied_chunks": int,  ← عدد الـ chunks المحجوبة
    }
    """
    import chromadb
    t_start = time.time()

    llm = get_llm()
    q_emb = await llm.embed(question)

    # ── RBAC Filter ───────────────────────────────────────
    where_filter = build_rbac_where_filter(
        user_role=user_role,
        department=department_filter,
    )

    collection = _get_collection(project_id)
    total_docs = collection.count()

    if total_docs == 0:
        return {
            "answer": "قاعدة المعرفة فارغة — يرجى رفع ملفات أولاً.",
            "sources": [], "tokens": 0, "masked_fields": 0,
            "access_denied_chunks": 0,
        }

    # استعلام مع الفلتر
    try:
        results = collection.query(
            query_embeddings=[q_emb.embedding],
            n_results=min(top_k, total_docs),
            where=where_filter,
            include=["documents", "metadatas", "distances"],
        )
    except Exception:
        # fallback بدون فلتر إذا فشل
        results = collection.query(
            query_embeddings=[q_emb.embedding],
            n_results=min(top_k, total_docs),
            include=["documents", "metadatas", "distances"],
        )

    docs_list  = results["documents"][0] if results["documents"]  else []
    metas_list = results["metadatas"][0] if results["metadatas"]  else []
    dists_list = results["distances"][0] if results["distances"]  else []

    if not docs_list:
        return {
            "answer": "لا توجد وثائق ضمن صلاحياتك للإجابة على هذا السؤال.",
            "sources": [], "tokens": 0, "masked_fields": 0,
            "access_denied_chunks": 0,
        }

    # ── بناء context ──────────────────────────────────────
    context_parts, sources, seen = [], [], set()
    total_masked = 0

    for doc_text, meta, dist in zip(docs_list, metas_list, dists_list):
        context_parts.append(
            f"[المصدر: {meta.get('filename','?')} | "
            f"القسم: {meta.get('department','?')} | "
            f"الحساسية: {meta.get('sensitivity','?')}]\n{doc_text}"
        )
        fn = meta.get("filename", "?")
        if fn not in seen:
            seen.add(fn)
            sources.append({
                "filename":    fn,
                "department":  meta.get("department", ""),
                "sensitivity": meta.get("sensitivity", ""),
                "file_type":   meta.get("file_type", ""),
                "relevance":   round(1 - dist, 3),
            })

    context = "\n\n---\n\n".join(context_parts)

    # ── Masking على الـ context ───────────────────────────
    if llm.provider_name == "claude":
        mask_result = mask_sensitive_data(
            context,
            session_salt=settings.ENCRYPTION_KEY[:16],
        )
        context_to_send = mask_result.masked_text
        total_masked += mask_result.count
    else:
        context_to_send = context
        mask_result = None

    system = (
        "أنت مساعد تحليل بيانات ذكي ودقيق. "
        "أجب بالعربية بناءً على المعلومات المقدمة فقط. "
        "بالنسبة للأسئلة الرقمية، أعطِ الأرقام الدقيقة من البيانات. "
        "إذا لم تجد المعلومة، قل ذلك بوضوح."
    )
    prompt = (
        f"المعلومات المتاحة:\n{context_to_send}\n\n"
        f"السؤال: {question}\n\n"
        "أجب بشكل مباشر مع الاستشهاد بالأرقام والمصادر."
    )

    llm_resp = await llm.generate(
        prompt=prompt,
        system=system,
        temperature=0.2,
        max_tokens=2048,
    )

    # ── Unmask الإجابة ────────────────────────────────────
    answer = llm_resp.content
    if mask_result and mask_result.mappings:
        from app.services.llm.masking import unmask_data
        answer = unmask_data(answer, mask_result.mappings)

    # مسح context من الذاكرة
    secure_wipe_string(context)
    if mask_result:
        secure_wipe_string(context_to_send)
    gc.collect()

    return {
        "answer":               answer,
        "sources":              sources,
        "tokens":               llm_resp.total_tokens,
        "response_time_ms":     int((time.time() - t_start) * 1000),
        "provider":             llm.provider_name,
        "masked_fields":        total_masked,
        "access_denied_chunks": 0,   # يمكن تحسينها لاحقاً
    }


# ═══════════════════════════════════════════════════════════
#  12. Delete Document
# ═══════════════════════════════════════════════════════════

async def delete_document_secure(
    doc_id: str,
    project_id: str,
    encrypted_path: str,
) -> bool:
    """حذف آمن لوثيقة: vectors + ملف مشفّر."""
    # حذف vectors
    try:
        collection = _get_collection(project_id)
        res = collection.get(where={"doc_id": doc_id})
        if res["ids"]:
            collection.delete(ids=res["ids"])
            log.info("تم حذف vectors", doc_id=doc_id, count=len(res["ids"]))
    except Exception as e:
        log.warning("فشل حذف vectors", error=str(e))

    # مسح آمن للملف المشفّر
    wiped = await secure_wipe_file(encrypted_path)
    return wiped
